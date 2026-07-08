#!/usr/bin/env python3
"""Generate a local transform smoke-test bundle from one Markdown article.

This script is intentionally provider-neutral: no vault path, personal profile,
API key, or external SaaS dependency is baked in. It exercises the public local
provider path: Markdown/report, PDF via browser print, PPTX, HTML deck, dense
infographic PNG, data table, quiz, flashcards, mindmap, podcast/video scripts.
"""

from __future__ import annotations

import argparse
import csv
import html
import json
import os
import re
import subprocess
import sys
import textwrap
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception as exc:  # pragma: no cover - reported at runtime
    Presentation = None
    PPTX_IMPORT_ERROR = exc
else:
    PPTX_IMPORT_ERROR = None


TERM_LIBRARY = [
    ("AI", "人工智能", "让机器完成原本需要人类智能参与的理解、生成、判断和执行任务。"),
    ("机器学习", "学习范式", "给机器大量样本，让它从数据里总结规律，而不是靠程序员写死规则。"),
    ("深度学习", "模型训练", "用多层神经网络逐级提取特征，是今天很多大模型的基础。"),
    ("LLM", "大语言模型", "理解和生成文本的模型，不是搜索引擎，也不是固定规则聊天机器人。"),
    ("GPT", "模型家族", "可以理解为 AI 能力发动机；ChatGPT 是面向人的产品入口。"),
    ("AIGC", "生成结果", "AI 生成的文字、图片、音频、视频、代码等内容产物。"),
    ("API", "系统接口", "让软件稳定调用 AI 能力，而不是让人打开聊天窗口手动输入。"),
    ("Prompt", "任务表达", "把角色、任务、对象、资料、格式、限制讲清楚的输入指令。"),
    ("System Prompt", "系统约束", "比普通输入更上层的角色、规则和边界说明。"),
    ("Token", "上下文单位", "模型处理文本时的基本切分单位，影响上下文容量和成本。"),
    ("上下文", "工作记忆", "模型当前能看到的材料、对话、规则和任务状态。"),
    ("知识库", "资料源", "把自己的资料交给 AI 参考，降低凭空编造和重复输入。"),
    ("RAG", "检索增强生成", "先从知识库查相关资料，再把资料交给模型生成回答。"),
    ("Embedding", "语义向量", "把文本转成可比较的向量，用于相似度检索。"),
    ("MCP", "工具协议", "让 AI 以标准方式连接外部工具、数据和服务。"),
    ("CLI", "命令行入口", "给 agent 或开发者调用工具的稳定命令接口。"),
    ("工具调用", "执行能力", "让 AI 从回答问题扩展到搜索、读写文件、调用服务。"),
    ("Agent", "自主执行体", "能理解目标、拆任务、调用工具、根据结果继续推进的 AI 系统。"),
    ("Workflow", "固定流程", "预先设计好的步骤链，稳定但自主性弱于 agent。"),
    ("Coding Agent", "编程代理", "能读代码、改代码、跑测试、修复问题的 agent。"),
    ("训练", "模型构建", "用大量数据和算力让模型学习通用能力。"),
    ("微调", "定向适配", "在已有模型上用领域数据继续训练，让它更适合某类任务。"),
    ("蒸馏", "模型压缩", "让小模型学习大模型输出，降低成本或部署门槛。"),
    ("量化", "推理优化", "降低参数精度以减少显存和计算开销。"),
    ("幻觉", "可靠性风险", "模型生成看似合理但没有依据或事实错误的内容。"),
    ("护栏", "安全边界", "用规则、校验和流程限制模型做不该做的事。"),
    ("Vibe Coding", "开发方式", "用自然语言和 agent 快速推动代码实现，但仍需要验证。"),
    ("GEO", "生成引擎优化", "面向 AI 搜索/回答环境组织内容可见性的方法。"),
]

MIN_DENSE_SLIDES = 14
MAX_DECK_SLIDES = 18
MIN_INFOGRAPHIC_BLOCKS = 12
MIN_REPORT_CHARS = 4500
MIN_PODCAST_CHARS = 1800

SECTION_THEMES = [
    ("能力底座", ("AI", "机器学习", "深度学习", "LLM", "GPT", "AIGC")),
    ("使用入口", ("模型", "产品入口", "API", "CLI")),
    ("需求表达", ("Prompt", "System Prompt", "提示词工程")),
    ("资料接入", ("Token", "上下文", "知识库", "RAG", "Embedding")),
    ("工具连接", ("MCP", "插件", "工具调用")),
    ("执行系统", ("Agent", "Coding Agent", "Workflow")),
    ("模型工程", ("训练", "微调", "蒸馏", "量化", "幻觉")),
    ("工程治理", ("上下文工程", "护栏", "Vibe Coding", "Harness", "GEO")),
]


@dataclass
class Article:
    path: Path
    title: str
    author: str
    url: str
    published_at: str
    body: str
    sections: list[tuple[str, str]]


def slugify(value: str) -> str:
    value = re.sub(r"[\\/:*?\"<>|#%{}^~\[\]`]+", "-", value).strip(" .-")
    return value[:90] or "article-transform"


def parse_frontmatter(text: str) -> tuple[dict[str, str], str]:
    if not text.startswith("---\n"):
        return {}, text
    end = text.find("\n---", 4)
    if end == -1:
        return {}, text
    raw = text[4:end].strip()
    rest = text[end + 4 :].lstrip()
    data: dict[str, str] = {}
    for line in raw.splitlines():
        if ":" not in line or line.startswith(" "):
            continue
        key, val = line.split(":", 1)
        data[key.strip()] = val.strip().strip('"')
    return data, rest


def strip_noise(line: str) -> str:
    line = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", line)
    line = line.replace("> [!source]- 来源信息", "")
    line = re.sub(r"^>\s?", "", line)
    line = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
    line = re.sub(r"`([^`]+)`", r"\1", line)
    return line.strip()


def parse_article(path: Path) -> Article:
    text = path.read_text(encoding="utf-8")
    fm, body = parse_frontmatter(text)
    title = fm.get("title") or ""
    m = re.search(r"^#\s+(.+)$", body, flags=re.MULTILINE)
    if m:
        title = m.group(1).strip()
    if not title:
        title = path.stem

    sections: list[tuple[str, str]] = []
    current_title = "导读"
    current_lines: list[str] = []
    for raw in body.splitlines():
        h = re.match(r"^##\s+(.+)$", raw)
        if h:
            if current_lines:
                sections.append((current_title, "\n".join(current_lines).strip()))
            current_title = h.group(1).strip()
            current_lines = []
            continue
        if raw.startswith("# "):
            continue
        cleaned = strip_noise(raw)
        if cleaned:
            current_lines.append(cleaned)
    if current_lines:
        sections.append((current_title, "\n".join(current_lines).strip()))

    return Article(
        path=path,
        title=title,
        author=fm.get("author", ""),
        url=fm.get("url", ""),
        published_at=fm.get("published_at", ""),
        body=body.strip(),
        sections=sections,
    )


def matched_terms(article: Article) -> list[tuple[str, str, str]]:
    haystack = article.body.lower()
    terms: list[tuple[str, str, str]] = []
    for term, category, definition in TERM_LIBRARY:
        if term.lower() in haystack or term in article.title:
            terms.append((term, category, definition))
    return terms or TERM_LIBRARY[:12]


def deck_slide_count(terms: list[tuple[str, str, str]], article: Article) -> int:
    """Choose a minimum slide count from source complexity, not from taste."""
    complexity = max(len(terms), len(article.sections) * 2)
    if complexity >= 20:
        return MAX_DECK_SLIDES
    if complexity >= 14:
        return 16
    return MIN_DENSE_SLIDES


def theme_rows(terms: list[tuple[str, str, str]]) -> list[tuple[str, list[tuple[str, str, str]]]]:
    names = {term: (term, category, definition) for term, category, definition in terms}
    rows: list[tuple[str, list[tuple[str, str, str]]]] = []
    for theme, wanted in SECTION_THEMES:
        bucket = [names[name] for name in wanted if name in names]
        if bucket:
            rows.append((theme, bucket))
    used = {term for _, bucket in rows for term, _, _ in bucket}
    leftovers = [item for item in terms if item[0] not in used]
    if leftovers:
        rows.append(("补充热词", leftovers))
    return rows


def qa_floor(article: Article, terms: list[tuple[str, str, str]]) -> dict[str, int]:
    slides = deck_slide_count(terms, article)
    return {
        "min_terms": min(len(terms), 20),
        "min_slides": slides,
        "min_infographic_blocks": max(MIN_INFOGRAPHIC_BLOCKS, min(len(terms), 18)),
        "min_report_chars": max(MIN_REPORT_CHARS, len(article.body) // 2),
        "min_podcast_chars": MIN_PODCAST_CHARS,
    }


def section_excerpt(text: str, limit: int = 170) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    return text[:limit].rstrip() + ("..." if len(text) > limit else "")


def write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content.strip() + "\n", encoding="utf-8")


def write_report(article: Article, out_dir: Path, terms: list[tuple[str, str, str]]) -> Path:
    toc = "\n".join(f"- {title}" for title, _ in article.sections[:14])
    coverage = "\n".join(
        f"| {theme} | {', '.join(term for term, _, _ in bucket)} | {len(bucket)} |"
        for theme, bucket in theme_rows(terms)
    )
    sections = "\n\n".join(
        f"### {title}\n\n{content.strip()}" for title, content in article.sections if content.strip()
    )
    glossary = "\n".join(f"| {t} | {c} | {d} |" for t, c, d in terms)
    checklist = "\n".join(
        f"- **{term}**：用于判断「{category}」这一层是否讲清楚；验收口径：{definition}"
        for term, category, definition in terms[:24]
    )
    report = f"""
# {article.title}｜保真转换报告

来源：{article.author or "unknown"} · {article.published_at or "unknown"} · {article.url or article.path}

## 转换说明

- 内容模式：`preserve + learning`
- 目标：保留原文章节、案例、术语和清单，再补充可学习、可检索、可导出的结构。
- 重要边界：本地 provider 不声称调用 NotebookLM；NotebookLM 产物由 notebooklm matrix 单独记录。

## 执行摘要

1. 这篇文章不是在解释单个 AI 产品，而是在用「内容生产小助手」串起 AI 工作系统：能力底座、使用入口、需求表达、资料接入、工具连接、执行系统和治理边界。
2. 转换时必须保留术语覆盖度。若 PPT、报告或音频只讲 AI / Prompt / RAG / Agent 四五个词，就遗漏了原文主体。
3. 文章的学习价值在「区分层级」：AI 与 AIGC、GPT 与 ChatGPT、RAG 与联网搜索、Agent 与 Workflow、提示词工程与上下文工程，都不能混成一页摘要。
4. 真正的落点是方法论：从一个需求出发，准备资料，表达任务，连接工具，执行交付，再用护栏和验收控制风险。

## 原文章节地图

{toc}

## 概念覆盖矩阵

| 模块 | 必须覆盖的术语 | 数量 |
|---|---|---:|
{coverage}

## 术语表

| 术语 | 类型 | 小白解释 |
|---|---|---|
{glossary}

## 学习与交付检查清单

{checklist}

## 保真正文

{sections}
"""
    path = out_dir / "report.md"
    write_text(path, report)
    return path


def write_data_table(out_dir: Path, terms: list[tuple[str, str, str]]) -> Path:
    path = out_dir / "data-table.csv"
    with path.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.writer(fh)
        writer.writerow(["term", "category", "plain_explanation", "learning_use"])
        for term, category, definition in terms:
            writer.writerow([term, category, definition, "用于概念辨析、试题、闪卡、PPT 知识地图"])
    return path


def write_quiz(out_dir: Path, terms: list[tuple[str, str, str]]) -> Path:
    selected = terms[:14]
    questions: list[str] = []
    answers: list[str] = []
    for idx, (term, category, definition) in enumerate(selected, 1):
        questions.append(
            f"### {idx}. {term} 属于什么概念？\n\n"
            f"A. {category}\nB. 一种图片格式\nC. 一种硬件接口\nD. 固定的营销口号\n"
        )
        answers.append(f"### {idx}. 答案：A\n\n解析：{definition}\n")
    next_id = len(selected) + 1
    questions.append(
        f"### {next_id}. 为什么把“转换文章为 PDF”直接做成几段摘要是不合格的？\n\n"
        "请用 3-5 句话回答。\n"
    )
    answers.append(
        f"### {next_id}. 答案要点\n\n用户要的是保真转换，不是总结；应保留章节、例子、关键清单和来源，只有明确说总结时才压缩。\n"
    )
    next_id += 1
    questions.append(
        f"### {next_id}. RAG 相比直接问 LLM 的核心改进是什么？\n\n"
        "请说明“先查资料，再回答”的价值。\n"
    )
    answers.append(
        f"### {next_id}. 答案要点\n\n先检索知识库中的相关材料，再让模型基于材料生成，减少凭空编造并增强可追溯性。\n"
    )
    next_id += 1
    questions.append(
        f"### {next_id}. 应用题：你要做一个“AI 内容生产小助手”，请按文章逻辑列出 5 个模块。\n\n"
        "至少包含 Prompt、RAG/MCP/Agent 或 Workflow 中的三个概念。\n"
    )
    answers.append(
        f"### {next_id}. 答案要点\n\n可包含：Prompt 表达任务、知识库/RAG 接资料、MCP/工具调用接外部能力、Agent 推进执行、Workflow 固定重复步骤、护栏/验收控制风险。\n"
    )
    path = out_dir / "quiz.md"
    write_text(path, "# 试卷：AI 名词入门自测\n\n## 题目\n\n" + "\n".join(questions) + "\n\n## 答案与解析\n\n" + "\n".join(answers))
    return path


def write_flashcards(out_dir: Path, terms: list[tuple[str, str, str]]) -> tuple[Path, Path]:
    md_lines = ["# 闪卡：AI 名词入门", "", "| 正面 | 背面 |", "|---|---|"]
    csv_path = out_dir / "flashcards.csv"
    with csv_path.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.writer(fh)
        writer.writerow(["front", "back"])
        for term, category, definition in terms:
            front = f"{term} 是什么？"
            back = f"{category}：{definition}"
            writer.writerow([front, back])
            md_lines.append(f"| {front} | {back} |")
    md_path = out_dir / "flashcards.md"
    write_text(md_path, "\n".join(md_lines))
    return md_path, csv_path


def write_mindmap(out_dir: Path, terms: list[tuple[str, str, str]]) -> Path:
    groups: dict[str, list[str]] = {}
    for term, category, _ in terms:
        groups.setdefault(category, []).append(term)
    lines = ["mindmap", "  root((AI 名词入门))"]
    for group, names in groups.items():
        lines.append(f"    {group}")
        for name in names[:6]:
            lines.append(f"      {name}")
    path = out_dir / "mindmap.mmd"
    write_text(path, "\n".join(lines))
    return path


def write_scripts(article: Article, out_dir: Path, terms: list[tuple[str, str, str]]) -> list[Path]:
    themes = theme_rows(terms)
    term_names = "、".join(t for t, _, _ in terms[:20])
    theme_script = "\n".join(
        f"### {idx}. {theme}\n"
        f"讲清：{', '.join(term for term, _, _ in bucket)}。\n"
        f"讲法：先给一句人话定义，再回到内容生产小助手案例，说明这一组概念在系统里解决什么问题。\n"
        for idx, (theme, bucket) in enumerate(themes, 1)
    )
    rehearsal_lines = "\n".join(
        f"- {theme}：听众应该能说出这一组概念的边界，并能把它放回“内容生产小助手”的链路中。"
        for theme, _ in themes
    )
    scene_lines = "\n\n".join(
        f"## 镜头 {idx + 1}｜{theme}\n"
        f"画面：{theme} 模块以流程板/关系图形式展开，关键节点依次亮起：{', '.join(term for term, _, _ in bucket[:5])}。\n"
        f"旁白：这一段解释 {theme}，重点不是背名词，而是知道它在 AI 工作系统里负责哪一环。"
        for idx, (theme, bucket) in enumerate(themes[:8], 1)
    )
    shot_lines = "\n".join(
        f"{idx + 3}. {theme}：镜头从具体任务推进到概念层级，画面中只出现短标签，不让中文小字承担密集解释。"
        for idx, (theme, _) in enumerate(themes[:10], 1)
    )
    podcast = f"""
# Podcast Script｜{article.title}

目标时长：8-10 分钟
目标听众：刚接触 AI 工具、但想把概念串成工作系统的小白

## 片头 0:00-0:40

今天我们不背术语表，而是用一个“AI 内容生产小助手”的案例，讲清楚这些词为什么会一起出现：{term_names}。
你手里有直播课逐字稿、学员问答、案例截图和产品资料，希望 AI 帮你整理成公众号文章、小红书文案和课程大纲。这个需求表面是一句 Prompt，拆开以后就是一套 AI 工作系统。

## 主体结构 0:40-7:40

{theme_script}

## 三组必须纠正的误解

1. LLM 不是搜索引擎。搜索找已有网页，LLM 生成回答；RAG 是把“先查资料”接到生成之前。
2. Prompt 不是魔法咒语。它至少要交代角色、任务、对象、资料、格式和限制。
3. Agent 不等于 Workflow。Workflow 是固定步骤链，Agent 会根据目标和反馈继续判断下一步。

## 主持人复述检查

这一段用于防止节目变成短摘要。录制前逐条确认：

{rehearsal_lines}

如果某一组讲不出来，就说明还没有覆盖原文主体，不能进入录制。

## 结尾 7:40-9:00

这篇文章最后真正要你记住的，不是某个热词，而是一条线：先分清 AI 能力底座，再选择产品/API/CLI 入口；用 Prompt 说清需求；用上下文、知识库、RAG 和 Embedding 接资料；用 MCP、插件、工具调用接外部世界；最后用 Agent、Coding Agent 或 Workflow 把事情推进到交付。只要涉及重要事实、钱、账号、数据安全，就必须加护栏和验收。
"""
    video = f"""
# Video Script｜教学长视频

目标时长：6-8 分钟
视频类型：白板讲解 + 概念地图 + 案例拆解

## 镜头 1｜开场任务
画面：课程主理人的资料堆：逐字稿、问答、案例图、产品文档。
旁白：你以为自己只是想“让 AI 帮我做内容”，其实已经进入 AI 工作系统。

{scene_lines}

## 镜头 {len(themes[:8]) + 2}｜总复盘
画面：一条横向链路：需求 -> Prompt -> 知识库/RAG -> MCP/工具 -> Agent/Workflow -> 产物 -> 验收。
旁白：真正有用的 AI，不只是回答，而是交付可验证结果。
"""
    cinematic = f"""
# Cinematic Video Shotlist｜概念电影感短片

目标：做 60-90 秒电影感概念片，用视觉隐喻承载情绪，用字幕/旁白承载事实；不要把密集中文塞进画面。

1. 冷开场：夜晚桌面，直播课逐字稿、学员问答、案例截图被屏幕光照亮。
2. 需求浮现：一句“做一个 AI 内容生产小助手”被拆成多张任务卡。
{shot_lines}
{len(themes[:10]) + 3}. 系统合流：Prompt、RAG、MCP、Agent、Workflow 汇入同一条工作流。
{len(themes[:10]) + 4}. 验收落点：文件、截图、来源链接和测试结果逐一盖章。
{len(themes[:10]) + 5}. 收束字幕：AI 是工作系统，不是孤立聊天框。
"""
    files = [
        (out_dir / "podcast-script.md", podcast),
        (out_dir / "video-script.md", video),
        (out_dir / "cinematic-video-shotlist.md", cinematic),
    ]
    for path, content in files:
        write_text(path, content)
    return [path for path, _ in files]


def html_doc(title: str, body: str, css: str) -> str:
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{html.escape(title)}</title>
<style>{css}</style>
</head>
<body>{body}</body>
</html>"""


def write_report_html(article: Article, out_dir: Path, terms: list[tuple[str, str, str]]) -> Path:
    term_rows = "".join(
        f"<tr><td>{html.escape(t)}</td><td>{html.escape(c)}</td><td>{html.escape(d)}</td></tr>"
        for t, c, d in terms
    )
    theme_cards = "".join(
        f"""
<article class="theme">
  <h3>{html.escape(theme)}</h3>
  <p>{html.escape(' / '.join(term for term, _, _ in bucket))}</p>
  <small>{len(bucket)} 个概念</small>
</article>
"""
        for theme, bucket in theme_rows(terms)
    )
    section_blocks = "".join(
        f"<section class='source-section'><h2>{html.escape(title)}</h2><p>{html.escape(section_excerpt(content, 1100))}</p></section>"
        for title, content in article.sections
        if content.strip()
    )
    css = """
*{box-sizing:border-box}body{margin:0;background:#e7edf0;color:#172026;font-family:"PingFang SC","Noto Sans CJK SC",sans-serif;line-height:1.58}
main{max-width:1160px;margin:0 auto;padding:42px;background:#f9f7f0}
.hero{background:#092734;color:#fff;border:3px solid #0a0f13;padding:36px 40px;box-shadow:12px 12px 0 #f0bd57}
h1{font-size:42px;line-height:1.12;margin:0 0 14px;letter-spacing:0}
.meta{font-size:14px;color:#bfd5db;margin-bottom:20px}
.verdict{font-size:24px;font-weight:800;color:#f7cf7a;margin-top:18px}
.summary{display:grid;grid-template-columns:repeat(4,1fr);gap:14px;margin:26px 0}
.summary div{background:#fff;border:2px solid #12232c;padding:16px;min-height:150px}
.summary b{display:block;font-size:18px;margin-bottom:8px;color:#0e4c68}.summary p{font-size:14.5px;margin:0}
h2{font-size:25px;margin:34px 0 12px;border-top:3px solid #102b36;padding-top:18px;color:#102b36}
.themes{display:grid;grid-template-columns:repeat(4,1fr);gap:12px}
.theme{background:#102b36;color:#f8f4e8;border:1px solid #5f8290;padding:14px;min-height:132px}
.theme h3{margin:0 0 8px;font-size:20px;color:#f7cf7a}.theme p{font-size:13.5px;margin:0 0 10px}.theme small{color:#b4d3dc}
p{font-size:15.5px;margin:0 0 12px}
table{border-collapse:collapse;width:100%;font-size:13.5px;margin:16px 0 28px;background:white}
td,th{border:1px solid #c7cdd1;padding:9px;vertical-align:top}
th{background:#102b36;color:#fff}
.source-section{background:#fff;border-left:7px solid #f0bd57;padding:16px 18px;margin:14px 0}
.source-section h2{font-size:20px;border:0;margin:0 0 8px;padding:0;color:#102b36}
.note{background:#d9eef8;border-left:6px solid #1976a3;padding:14px;margin:22px 0}
.risk{background:#fae5e5;border-left:6px solid #bf3e48;padding:14px;margin:18px 0}
@media print{body{background:white}main{max-width:none;padding:22px}.hero{box-shadow:none}h1{font-size:30px}.summary,.themes{grid-template-columns:repeat(2,1fr)}}
"""
    body = f"""
<main>
  <section class="hero">
    <h1>{html.escape(article.title)}</h1>
    <div class="meta">视觉报告 · {html.escape(article.author or "unknown")} · {html.escape(article.published_at or "")}</div>
    <div class="verdict">主判断：AI 不是一个孤立聊天框，而是从理解需求到连接资料、调用工具、执行交付的工作系统。</div>
  </section>
  <section class="summary">
    <div><b>不是摘要</b><p>报告必须覆盖原文十几个概念和章节链路，不能只保留 3-5 个观点。</p></div>
    <div><b>先分层</b><p>能力底座、使用入口、资料接入、工具连接、执行系统分别承担不同角色。</p></div>
    <div><b>再串线</b><p>内容生产小助手案例把 Prompt、RAG、MCP、Agent 和 Workflow 串成执行链。</p></div>
    <div><b>最后验收</b><p>幻觉、上下文缺失、工具失败都要靠来源、测试和护栏兜住。</p></div>
  </section>
  <h2>术语表</h2>
  <table><thead><tr><th>术语</th><th>类型</th><th>小白解释</th></tr></thead><tbody>{term_rows}</tbody></table>
  <h2>概念覆盖地图</h2>
  <section class="themes">{theme_cards}</section>
  <div class="risk">验收红线：如果 PPT、报告、播客或试卷没有覆盖这张地图的大部分模块，就只是摘要，不算转换完成。</div>
  <h2>原文章节与证据</h2>
  {section_blocks}
</main>
"""
    path = out_dir / "report.html"
    write_text(path, html_doc(article.title, body, css))
    return path


def write_infographic_html(article: Article, out_dir: Path, terms: list[tuple[str, str, str]]) -> Path:
    cards = "".join(
        f"""
<div class="card c{idx % 5}">
  <b>{idx:02d} {html.escape(term)}</b>
  <span>{html.escape(definition)}</span>
</div>
"""
        for idx, (term, _, definition) in enumerate(terms[:18], 1)
    )
    flow = [
        ("需求", "把资料变成内容"),
        ("Prompt", "说清角色、任务、对象、格式"),
        ("RAG", "从知识库取证据"),
        ("MCP", "连接外部工具"),
        ("Agent", "推进执行并验收"),
        ("交付", "文章、卡片、课程、发布"),
    ]
    flow_html = "".join(f"<li><b>{html.escape(a)}</b><span>{html.escape(b)}</span></li>" for a, b in flow)
    css = """
*{box-sizing:border-box}body{margin:0;background:#061923;color:#f8f4e8;font-family:"PingFang SC","Noto Sans CJK SC",sans-serif}
.poster{width:1080px;height:1920px;margin:0 auto;padding:32px;background:#061923;display:flex;flex-direction:column;gap:12px}
.top{border:2px solid #e4b363;padding:22px;background:#102b36}
h1{font-size:44px;line-height:1.14;margin:0 0 8px}.sub{font-size:21px;color:#cce8e2}
.verdict{margin-top:14px;padding:16px;background:#813c2a;border:1px solid #ffca6e;font-size:34px;font-weight:800;text-align:center}
.grid{display:grid;grid-template-columns:repeat(3,1fr);gap:10px}
.card{min-height:116px;border:1px solid #4d7180;background:#0d2632;padding:12px;border-radius:8px}
.card b{display:block;font-size:22px;color:#ffdc84;margin-bottom:6px}.card span{font-size:16.5px;line-height:1.32;color:#e9f5f1}
.c1 b{color:#6fb6ff}.c2 b{color:#ff7777}.c3 b{color:#87d37c}.c4 b{color:#f6c667}
.two{display:grid;grid-template-columns:1.05fr .95fr;gap:14px}
.panel{border:1px solid #567989;border-radius:8px;padding:16px;background:#0b2330}
.panel h2{font-size:28px;margin:0 0 10px;color:#ffdc84}
.flow{list-style:none;margin:0;padding:0;display:grid;grid-template-columns:repeat(2,1fr);gap:10px}
.flow li{padding:11px;border-left:6px solid #6fb6ff;background:#102f3e}.flow b{font-size:21px;display:block}.flow span{font-size:16px}
.compare{display:grid;grid-template-columns:1fr 1fr;gap:10px}.compare div{padding:14px;border-radius:8px;font-size:18px;line-height:1.42}
.yes{background:#123d2b;border:1px solid #47bf72}.risk{background:#4a171e;border:1px solid #ff6b7a}
.strip{display:grid;grid-template-columns:repeat(4,1fr);gap:10px}.strip div{border:1px solid #4d7180;background:#102b36;border-radius:8px;padding:12px;min-height:118px}.strip b{font-size:20px;color:#ffdc84}.strip p{font-size:15.5px;line-height:1.34;margin:6px 0 0}
.bottom{display:grid;grid-template-columns:1fr 1fr 1fr;gap:12px}.bottom .panel{min-height:210px}.bottom p{font-size:17px;line-height:1.45}
.foot{font-size:15px;color:#a8c3c8;text-align:center}
"""
    body = f"""
<main class="poster">
  <section class="top">
    <h1>{html.escape(article.title)}</h1>
    <div class="sub">一个案例看懂：模型、资料、工具、流程如何组成 AI 工作系统</div>
    <div class="verdict">不是孤立聊天框，而是可验证的工作流</div>
  </section>
  <section class="grid">{cards}</section>
  <section class="two">
    <div class="panel"><h2>内容生产小助手流程</h2><ul class="flow">{flow_html}</ul></div>
    <div class="panel"><h2>三组易混概念</h2>
      <div class="compare">
        <div class="yes"><b>AI / AIGC</b><br>AI 是能力，AIGC 是生成结果。</div>
        <div class="risk"><b>LLM / 搜索</b><br>搜索找网页，LLM 生成回答，RAG 把二者接起来。</div>
        <div class="yes"><b>Agent / Workflow</b><br>Agent 会根据反馈推进；Workflow 是固定步骤链。</div>
        <div class="risk"><b>Prompt / 上下文</b><br>Prompt 是任务表达，上下文是模型当前能看到的工作记忆。</div>
      </div>
    </div>
  </section>
  <section class="strip">
    <div><b>Prompt 六要素</b><p>角色、任务、对象、资料、格式、限制；缺项就会把自由度交给模型。</p></div>
    <div><b>RAG 三步</b><p>切分资料 -> 向量检索 -> 带证据生成；关键是来源可追溯。</p></div>
    <div><b>Agent 判断</b><p>是否会拆任务、调用工具、读取反馈并继续推进，而不是只回复一句话。</p></div>
    <div><b>验收证据</b><p>输出文件、来源链接、测试结果、人工审批点都要留下。</p></div>
  </section>
  <section class="bottom">
    <div class="panel"><h2>支持使用</h2><p>用清楚的角色、资料、格式、限制来控制输出。</p><p>把知识库和工具接入，减少重复劳动。</p></div>
    <div class="panel"><h2>必须警惕</h2><p>没有来源的答案可能是幻觉。</p><p>Agent 做完不等于正确，必须验收。</p></div>
    <div class="panel"><h2>继续追问</h2><p>我的资料放在哪里？</p><p>哪些步骤可自动化？哪些必须人工审批？</p></div>
  </section>
  <div class="foot">Local visual provider smoke · source: {html.escape(article.url or str(article.path))}</div>
</main>
"""
    path = out_dir / "infographic.html"
    write_text(path, html_doc(article.title, body, css))
    return path


def slide(title: str, body: str, kind: str = "text") -> str:
    return f'<section class="slide {kind}"><div class="paper">{body}</div><footer>{html.escape(title)}</footer></section>'


def write_deck_html(article: Article, out_dir: Path, terms: list[tuple[str, str, str]]) -> Path:
    themes = theme_rows(terms)
    floor = qa_floor(article, terms)
    term_cards = "".join(
        f"<div><b>{html.escape(t)}</b><span>{html.escape(d)}</span></div>" for t, _, d in terms[:12]
    )
    theme_cards = "".join(
        f"<div><b>{html.escape(theme)}</b><span>{html.escape(' / '.join(term for term, _, _ in bucket))}</span></div>"
        for theme, bucket in themes[:8]
    )
    glossary_rows_a = "".join(
        f"<tr><td>{html.escape(t)}</td><td>{html.escape(c)}</td><td>{html.escape(d)}</td></tr>"
        for t, c, d in terms[:9]
    )
    glossary_rows_b = "".join(
        f"<tr><td>{html.escape(t)}</td><td>{html.escape(c)}</td><td>{html.escape(d)}</td></tr>"
        for t, c, d in terms[9:18]
    )
    css = """
*{box-sizing:border-box}body{margin:0;background:#2d2d2d;font-family:"PingFang SC","Noto Sans CJK SC",sans-serif;color:#171717}
.deck{width:1920px;margin:0 auto}.slide{width:1920px;height:1080px;position:relative;padding:72px;background:#2d2d2d;page-break-after:always}
.paper{position:absolute;inset:70px 145px 76px 92px;background:#f8f6f1;border:2px solid #252525;padding:64px 78px;box-shadow:18px 18px 0 #98d4bb}
.slide:nth-child(2n) .paper{box-shadow:18px 18px 0 #c7b8ea}.slide:nth-child(3n) .paper{box-shadow:18px 18px 0 #f4b8c5}
h1{font-size:64px;line-height:1.08;margin:0 0 22px;font-weight:850}h2{font-size:54px;line-height:1.08;margin:0 0 30px}
p,li{font-size:28px;line-height:1.42}.kicker{font-size:22px;color:#6b6358;margin-bottom:28px}.grid{display:grid;grid-template-columns:repeat(4,1fr);gap:18px}
.grid div,.box{border:2px solid #252525;padding:20px;background:#fffdf8;min-height:132px}.grid b{display:block;font-size:27px;margin-bottom:8px}.grid span{font-size:19px;line-height:1.35}
.flow{display:grid;grid-template-columns:repeat(6,1fr);gap:14px}.flow div{border-top:8px solid #98d4bb;padding:16px;background:#fffdf8;min-height:220px}
.flow b{font-size:30px;display:block}.two{display:grid;grid-template-columns:1fr 1fr;gap:24px}.risk{background:#ffd9df}.ok{background:#ddf4e8}
.table{width:100%;border-collapse:collapse;font-size:22px}.table td,.table th{border:2px solid #252525;padding:12px;text-align:left;vertical-align:top}.table th{background:#102b36;color:#fff}
.matrix{display:grid;grid-template-columns:repeat(3,1fr);gap:18px}.matrix div{border:2px solid #252525;background:#fffdf8;padding:22px;min-height:190px}.matrix b{font-size:30px;display:block;margin-bottom:10px}.matrix p{font-size:23px;margin:0}
.cover-map{display:grid;grid-template-columns:repeat(6,1fr);gap:12px;margin-top:74px}.cover-map div{border:2px solid #252525;background:#fffdf8;padding:16px;min-height:128px}.cover-map b{display:block;font-size:25px;margin-bottom:8px}.cover-map span{font-size:18px;line-height:1.35}
footer{position:absolute;left:120px;bottom:34px;color:#f8f6f1;font-size:22px}.num{position:absolute;right:170px;top:90px;font-size:110px;color:#d8d2c4;opacity:.55}
@media print{body{background:white}.deck{width:auto}.slide{margin:0;break-after:page}}
"""
    slides = [
        slide(article.title, f"<div class='kicker'>AI 名词入门 · {floor['min_slides']} 页课程型 deck</div><h1>{html.escape(article.title)}</h1><p>用“AI 内容生产小助手”这个案例，把模型、资料、工具、执行和治理串成一张完整地图。</p><div class='cover-map'><div><b>能力</b><span>AI / LLM / AIGC</span></div><div><b>入口</b><span>产品 / API / CLI</span></div><div><b>表达</b><span>Prompt / System Prompt</span></div><div><b>资料</b><span>Token / RAG / Embedding</span></div><div><b>工具</b><span>MCP / Tool Calling</span></div><div><b>执行</b><span>Agent / Workflow</span></div></div>", "cover"),
        slide("文章地图", f"<h2>不是背词，而是搭系统</h2><div class='grid'>{theme_cards}</div>"),
        slide("案例拆解", "<h2>一句需求背后有 6 个模块</h2><div class='flow'><div><b>资料</b><p>逐字稿 / 问答 / 截图</p></div><div><b>表达</b><p>Prompt / 约束</p></div><div><b>检索</b><p>知识库 / RAG</p></div><div><b>工具</b><p>MCP / 插件</p></div><div><b>执行</b><p>Agent / Workflow</p></div><div><b>验收</b><p>来源 / 测试 / 护栏</p></div></div>"),
        slide("能力底座", "<h2>先分清 AI、机器学习、深度学习</h2><div class='matrix'><div><b>AI</b><p>让机器完成理解、生成、判断和执行。</p></div><div><b>机器学习</b><p>从样本中学习规律，不靠写死规则。</p></div><div><b>深度学习</b><p>用多层神经网络提取特征，是大模型基础。</p></div></div>"),
        slide("LLM / GPT / AIGC", "<h2>能力、模型、结果不是一回事</h2><table class='table'><tr><th>词</th><th>人话解释</th><th>常见误区</th></tr><tr><td>LLM</td><td>生成语言回答的大模型</td><td>不是搜索引擎</td></tr><tr><td>GPT</td><td>模型家族/能力发动机</td><td>不等于 ChatGPT 产品</td></tr><tr><td>AIGC</td><td>AI 生成的内容结果</td><td>不是某个单独工具</td></tr></table>"),
        slide("产品入口与 API", "<h2>人用产品，软件用 API</h2><div class='two'><div class='box'><b>ChatGPT / Kimi / 豆包</b><p>适合人打开窗口对话。</p></div><div class='box'><b>API / CLI</b><p>适合工具稳定调用 AI 能力，接进工作流。</p></div></div>"),
        slide("Prompt 六要素", "<h2>Prompt 不是一句口令</h2><div class='grid'><div><b>角色</b><span>以什么身份工作</span></div><div><b>任务</b><span>要完成什么</span></div><div><b>对象</b><span>处理哪份资料</span></div><div><b>格式</b><span>输出成什么</span></div><div><b>资料</b><span>依据在哪里</span></div><div><b>限制</b><span>不能做什么</span></div></div>"),
        slide("Token / 上下文", "<h2>AI 当前能看到什么，决定它能基于什么回答</h2><div class='two'><div class='box'><b>Token</b><p>模型处理文本的基本单位，影响成本、速度和上下文容量。</p></div><div class='box'><b>上下文</b><p>模型当前看到的材料、规则、对话和任务状态。</p></div></div>"),
        slide("知识库 / RAG / Embedding", "<h2>RAG 是先翻资料，再回答</h2><div class='flow'><div><b>资料</b><p>文档/笔记/问答</p></div><div><b>切分</b><p>把材料拆块</p></div><div><b>Embedding</b><p>转成语义向量</p></div><div><b>检索</b><p>找相关片段</p></div><div><b>生成</b><p>带证据回答</p></div><div><b>引用</b><p>可追溯</p></div></div>"),
        slide("MCP / 插件 / 工具调用", "<h2>让 AI 从会说话变成会干活</h2><p>MCP 像 AI 的 USB-C：工具方按统一接口提供能力，AI 就能更容易接上文档、浏览器、数据库、搜索、图片生成和发布工具。</p>"),
        slide("Agent / Coding Agent / Workflow", "<h2>从回答到执行，有三种层级</h2><table class='table'><tr><th>类型</th><th>特点</th><th>适合任务</th></tr><tr><td>普通聊天</td><td>给建议和文本</td><td>问答、解释</td></tr><tr><td>Workflow</td><td>固定步骤链</td><td>重复流程</td></tr><tr><td>Agent</td><td>会拆任务、调工具、看反馈</td><td>开放式交付</td></tr></table>"),
        slide("训练 / 微调 / 蒸馏 / 量化", "<h2>这些词解释模型怎么被造出来</h2><div class='grid'><div><b>训练</b><span>学习通用能力</span></div><div><b>微调</b><span>适配特定领域</span></div><div><b>蒸馏</b><span>让小模型学习大模型</span></div><div><b>量化</b><span>降低精度换部署效率</span></div></div>"),
        slide("幻觉 / 护栏 / 上下文工程", "<h2>会用 AI，是会准备环境和边界</h2><div class='two'><div class='risk'><p>幻觉：看似合理但没有依据。不能靠语气判断真假。</p></div><div class='ok'><p>护栏：规则、校验、审批、测试和来源引用。</p></div></div>"),
        slide("2026 热词", "<h2>Vibe Coding、Harness、GEO 都在回答同一个问题</h2><p>AI 能干活以后，人怎么指挥好、管得住、让产物被发现和复用。</p>"),
        slide("术语速查表 A", f"<h2>覆盖清单 A：能力、入口、表达</h2><table class='table'><tr><th>术语</th><th>类型</th><th>解释</th></tr>{glossary_rows_a}</table>"),
        slide("术语速查表 B", f"<h2>覆盖清单 B：资料、工具、执行</h2><table class='table'><tr><th>术语</th><th>类型</th><th>解释</th></tr>{glossary_rows_b}</table>"),
        slide("易混概念", "<h2>四组最容易混的判断</h2><div class='matrix'><div><b>AI vs AIGC</b><p>能力 vs 结果</p></div><div><b>GPT vs ChatGPT</b><p>模型家族 vs 产品入口</p></div><div><b>RAG vs 搜索</b><p>私有资料检索 vs 外部网页查找</p></div><div><b>Agent vs Workflow</b><p>自主推进 vs 固定步骤</p></div><div><b>Prompt vs 上下文</b><p>任务表达 vs 工作记忆</p></div><div><b>完成 vs 正确</b><p>产物存在 vs 经验证据</p></div></div>"),
        slide("行动清单", "<h2>读完就能做的 7 步</h2><ol><li>把任务写成 Prompt 六要素。</li><li>把资料放进知识库。</li><li>判断是否需要 RAG。</li><li>列出需要调用的工具。</li><li>区分 Agent 与 Workflow。</li><li>给每一步设计验收证据。</li><li>把产物链接回 Obsidian。</li></ol>"),
        slide("自测", "<h2>三道检查题</h2><ol><li>为什么 LLM 不等于搜索引擎？</li><li>什么时候用 Workflow，什么时候用 Agent？</li><li>如果 AI 生成结果很顺，但没有来源，应该怎么验收？</li></ol>"),
        slide("来源", f"<h2>Source</h2><p>{html.escape(article.author or 'unknown')}</p><p>{html.escape(article.url or str(article.path))}</p>"),
    ]
    body = f"<main class='deck'>{''.join(slides)}</main>"
    path = out_dir / "deck.html"
    write_text(path, html_doc(article.title, body, css))
    return path


def add_textbox(slide_obj, left, top, width, height, text, size=24, bold=False, color=(28, 28, 28)):
    box = slide_obj.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def write_deck_pptx(article: Article, out_dir: Path, terms: list[tuple[str, str, str]]) -> Path | None:
    if Presentation is None:
        print(f"python-pptx unavailable: {PPTX_IMPORT_ERROR}", file=sys.stderr)
        return None
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    bg = RGBColor(248, 246, 241)
    ink = RGBColor(26, 26, 26)
    mint = RGBColor(152, 212, 187)
    lavender = RGBColor(199, 184, 234)
    pink = RGBColor(244, 184, 197)

    def base_slide(title: str, accent=mint):
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
        shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.2), Inches(0.2), Inches(12.93), Inches(7.1))
        shape.fill.background()
        shape.line.color.rgb = ink
        shape.line.width = Pt(1.5)
        bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.2), Inches(6.9), Inches(12.93), Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        add_textbox(s, Inches(0.55), Inches(0.42), Inches(11.8), Inches(0.5), title, 18, True)
        return s

    s = base_slide("AI 名词入门", lavender)
    add_textbox(s, Inches(0.75), Inches(1.25), Inches(11.4), Inches(1.5), article.title, 38, True)
    add_textbox(s, Inches(0.78), Inches(3.25), Inches(8.8), Inches(1.1), "用一个内容生产助手案例，把模型、资料、工具、流程和执行边界一次讲清楚。", 23)
    cover_items = [
        ("能力", "AI / LLM / AIGC"),
        ("入口", "产品 / API / CLI"),
        ("表达", "Prompt / System Prompt"),
        ("资料", "Token / RAG / Embedding"),
        ("工具", "MCP / Tool Calling"),
        ("执行", "Agent / Workflow"),
    ]
    for idx, (head, text) in enumerate(cover_items):
        left = Inches(0.78 + (idx % 3) * 4.05)
        top = Inches(4.35 + (idx // 3) * 0.7)
        shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(3.75), Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 253, 248)
        shape.line.color.rgb = ink
        add_textbox(s, left + Inches(0.1), top + Inches(0.07), Inches(3.5), Inches(0.22), head, 12, True)
        add_textbox(s, left + Inches(0.1), top + Inches(0.29), Inches(3.5), Inches(0.22), text, 10)
    add_textbox(s, Inches(0.78), Inches(6.0), Inches(10.5), Inches(0.5), f"Source: {article.author or 'unknown'} · {article.url or article.path}", 12)

    pages = [
        ("文章地图", [f"{theme}：{', '.join(term for term, _, _ in bucket)}" for theme, bucket in theme_rows(terms)[:6]]),
        ("案例拆解", ["输入：逐字稿、问答、截图、产品资料", "输出：公众号、小红书、课程大纲、配图和发布清单", "链路：需求 -> 资料 -> 工具 -> 执行 -> 验收"]),
        ("能力底座", ["AI：完成理解、生成、判断和执行", "机器学习：从样本学习规律", "深度学习：多层神经网络提取特征"]),
        ("LLM / GPT / AIGC", ["LLM 生成语言回答，不等于搜索引擎", "GPT 是模型家族/能力发动机", "AIGC 是 AI 生成的内容结果"]),
        ("使用入口", ["产品入口给人使用", "API 给软件系统稳定调用", "CLI 给 agent / 开发者自动化调用"]),
        ("Prompt 六要素", ["角色", "任务", "对象", "资料", "格式", "限制"]),
        ("Token / 上下文", ["Token 影响成本、速度和容量", "上下文决定模型当前能基于什么回答", "资料越多，越需要结构化管理"]),
        ("知识库 / RAG / Embedding", ["知识库：自己的资料源", "Embedding：文本语义向量", "RAG：先检索资料，再生成回答"]),
        ("MCP / 插件 / 工具调用", ["MCP 标准化连接外部工具", "插件和工具调用让 AI 获取新能力", "联网搜索本质也是工具连接"]),
        ("Agent / Workflow", ["Workflow：固定步骤链，稳定可复用", "Agent：目标驱动，会根据反馈推进", "实际工作经常组合使用"]),
        ("Coding Agent", ["不只生成代码，还能读、写、跑、改、交付", "代表工具：Claude Code、Codex、Cursor 等", "责任仍然需要人验收"]),
        ("训练 / 微调 / 蒸馏 / 量化", ["训练：构建通用能力", "微调：适配场景", "蒸馏/量化：降低成本和部署门槛"]),
        ("幻觉与护栏", ["幻觉是大模型生成方式的天然风险", "护栏靠规则、来源、测试和审批", "语气顺不等于事实真"]),
        ("上下文工程", ["不是只写 Prompt", "还要准备代码、资料、规则、范例和验收标准", "CLAUDE.md / Skills / Harness 都是在准备边界"]),
        ("2026 热词", ["Vibe Coding：自然语言推动代码实现", "Harness：用流程和规则管住 agent", "GEO：面向 AI 生成环境组织内容可见性"]),
        ("术语覆盖检查", [f"{t}：{d}" for t, _, d in terms[:6]]),
        ("行动清单", ["把任务写成六要素", "把资料沉淀为知识库", "判断是否需要 RAG", "区分 Agent 与 Workflow", "给每一步设计验收证据"]),
        ("来源与边界", [f"作者：{article.author or 'unknown'}", f"来源：{article.url or article.path}", "本地 PPTX 是模板化 smoke，不声称调用 NotebookLM/Open Design"]),
    ]
    accents = [mint, lavender, pink]
    for i, (title, bullets) in enumerate(pages, 1):
        s = base_slide(title, accents[i % 3])
        add_textbox(s, Inches(0.75), Inches(1.0), Inches(4.2), Inches(0.8), f"0{i+1}", 44, True, (107, 99, 88))
        top = 1.75
        for b in bullets[:6]:
            add_textbox(s, Inches(1.0), Inches(top), Inches(11.0), Inches(0.55), f"- {b}", 21)
            top += 0.72

    path = out_dir / "deck.pptx"
    prs.save(path)
    return path


def render_with_playwright(out_dir: Path, node_bin: str, node_path: str | None) -> dict[str, str]:
    render_script = out_dir / ".render.cjs"
    render_script.write_text(
        r"""
const { chromium } = require('playwright');
const path = require('path');
const { pathToFileURL } = require('url');

function chromeLaunchOptions() {
  const candidates = [
    process.env.PLAYWRIGHT_CHROMIUM_EXECUTABLE,
    '/Applications/Google Chrome.app/Contents/MacOS/Google Chrome',
    '/Applications/Chromium.app/Contents/MacOS/Chromium'
  ].filter(Boolean);
  for (const executablePath of candidates) {
    try {
      require('fs').accessSync(executablePath);
      return { headless: true, executablePath };
    } catch (_) {}
  }
  return { headless: true };
}

async function renderHtml(input, opts) {
  const browser = await chromium.launch(chromeLaunchOptions());
  const page = await browser.newPage({ viewport: { width: opts.width, height: opts.height }, deviceScaleFactor: opts.scale || 1 });
  await page.goto(pathToFileURL(input).href, { waitUntil: 'networkidle' });
  if (opts.screenshot) {
    await page.screenshot({ path: opts.screenshot, fullPage: Boolean(opts.fullPage) });
  }
  if (opts.pdf) {
    await page.pdf({ path: opts.pdf, printBackground: true, format: opts.format || 'A4', margin: opts.margin || { top: '12mm', bottom: '12mm', left: '10mm', right: '10mm' } });
  }
  await browser.close();
}

(async () => {
  const root = process.argv[2];
  await renderHtml(path.join(root, 'infographic.html'), { width: 1080, height: 1920, screenshot: path.join(root, 'infographic.png') });
  await renderHtml(path.join(root, 'deck.html'), { width: 1920, height: 1080, screenshot: path.join(root, 'deck-cover.png') });
  await renderHtml(path.join(root, 'report.html'), { width: 1240, height: 1754, pdf: path.join(root, 'report.pdf'), fullPage: true });
})();
""",
        encoding="utf-8",
    )
    env = os.environ.copy()
    if node_path:
        env["NODE_PATH"] = node_path
    proc = subprocess.run(
        [node_bin, str(render_script), str(out_dir)],
        cwd=str(out_dir),
        env=env,
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )
    return {"returncode": str(proc.returncode), "stdout": proc.stdout, "stderr": proc.stderr}


def count_markdown_questions(text: str) -> tuple[int, int]:
    if "## 答案与解析" not in text:
        return (len(re.findall(r"^###\s+\d+\.", text, flags=re.MULTILINE)), 0)
    question_text, answer_text = text.split("## 答案与解析", 1)
    q_count = len(re.findall(r"^###\s+\d+\.", question_text, flags=re.MULTILINE))
    a_count = len(re.findall(r"^###\s+\d+\.", answer_text, flags=re.MULTILINE))
    return q_count, a_count


def validate_bundle(article: Article, out_dir: Path, terms: list[tuple[str, str, str]]) -> dict[str, object]:
    floor = qa_floor(article, terms)
    checks: dict[str, object] = {"thresholds": floor, "items": {}, "ok": True}

    def item(name: str, ok: bool, **extra: object) -> None:
        checks["items"][name] = {"ok": ok, **extra}
        if not ok:
            checks["ok"] = False

    report_md = out_dir / "report.md"
    report_text = report_md.read_text(encoding="utf-8") if report_md.exists() else ""
    covered_terms = [term for term, _, _ in terms if term in report_text]
    item(
        "report_md",
        report_md.exists() and len(report_text) >= floor["min_report_chars"] and len(covered_terms) >= floor["min_terms"],
        chars=len(report_text),
        covered_terms=len(covered_terms),
        min_chars=floor["min_report_chars"],
        min_terms=floor["min_terms"],
    )

    deck_html = out_dir / "deck.html"
    deck_text = deck_html.read_text(encoding="utf-8") if deck_html.exists() else ""
    html_slides = len(re.findall(r"class=['\"]slide", deck_text))
    item("deck_html", deck_html.exists() and html_slides >= floor["min_slides"], slides=html_slides, min_slides=floor["min_slides"])

    pptx = out_dir / "deck.pptx"
    ppt_slides = 0
    if pptx.exists() and Presentation is not None:
        try:
            ppt_slides = len(Presentation(str(pptx)).slides)
        except Exception:  # noqa: BLE001
            ppt_slides = 0
    item("deck_pptx", pptx.exists() and ppt_slides >= floor["min_slides"], slides=ppt_slides, min_slides=floor["min_slides"])

    infographic = out_dir / "infographic.html"
    info_text = infographic.read_text(encoding="utf-8") if infographic.exists() else ""
    info_blocks = len(re.findall(r'class="card ', info_text)) + len(re.findall(r'<li><b>', info_text)) + len(re.findall(r'<div><b>', info_text))
    item(
        "infographic_html",
        infographic.exists() and info_blocks >= floor["min_infographic_blocks"],
        blocks=info_blocks,
        min_blocks=floor["min_infographic_blocks"],
    )

    quiz = out_dir / "quiz.md"
    quiz_text = quiz.read_text(encoding="utf-8") if quiz.exists() else ""
    q_count, a_count = count_markdown_questions(quiz_text)
    item("quiz", quiz.exists() and q_count >= 12 and q_count == a_count, questions=q_count, answers=a_count)

    flashcards = out_dir / "flashcards.csv"
    flash_count = 0
    if flashcards.exists():
        with flashcards.open(encoding="utf-8", newline="") as fh:
            flash_count = max(0, sum(1 for _ in csv.reader(fh)) - 1)
    item("flashcards", flashcards.exists() and flash_count >= floor["min_terms"], cards=flash_count, min_cards=floor["min_terms"])

    data_table = out_dir / "data-table.csv"
    data_rows = 0
    if data_table.exists():
        with data_table.open(encoding="utf-8", newline="") as fh:
            data_rows = max(0, sum(1 for _ in csv.reader(fh)) - 1)
    item("data_table", data_table.exists() and data_rows >= floor["min_terms"], rows=data_rows, min_rows=floor["min_terms"])

    podcast = out_dir / "podcast-script.md"
    podcast_text = podcast.read_text(encoding="utf-8") if podcast.exists() else ""
    item("podcast_script", podcast.exists() and len(podcast_text) >= floor["min_podcast_chars"], chars=len(podcast_text), min_chars=floor["min_podcast_chars"])

    video = out_dir / "video-script.md"
    video_text = video.read_text(encoding="utf-8") if video.exists() else ""
    item("video_script", video.exists() and len(re.findall(r"^## 镜头", video_text, flags=re.MULTILINE)) >= 8, scenes=len(re.findall(r"^## 镜头", video_text, flags=re.MULTILINE)))

    cinematic = out_dir / "cinematic-video-shotlist.md"
    cinematic_text = cinematic.read_text(encoding="utf-8") if cinematic.exists() else ""
    item("cinematic_shotlist", cinematic.exists() and len(re.findall(r"^\d+\.", cinematic_text, flags=re.MULTILINE)) >= 10, shots=len(re.findall(r"^\d+\.", cinematic_text, flags=re.MULTILINE)))

    return checks


def build_bundle(args: argparse.Namespace) -> dict[str, object]:
    article = parse_article(Path(args.article).expanduser().resolve())
    out_dir = Path(args.out_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    terms = matched_terms(article)

    files: list[Path] = []
    files.append(write_report(article, out_dir, terms))
    files.append(write_report_html(article, out_dir, terms))
    files.append(write_data_table(out_dir, terms))
    files.append(write_quiz(out_dir, terms))
    files.extend(write_flashcards(out_dir, terms))
    files.append(write_mindmap(out_dir, terms))
    files.extend(write_scripts(article, out_dir, terms))
    files.append(write_infographic_html(article, out_dir, terms))
    files.append(write_deck_html(article, out_dir, terms))
    pptx_path = write_deck_pptx(article, out_dir, terms)
    if pptx_path:
        files.append(pptx_path)

    render_result = None
    if not args.no_render:
        render_result = render_with_playwright(out_dir, args.node_bin, args.node_path)
        for name in ["report.pdf", "infographic.png", "deck-cover.png"]:
            p = out_dir / name
            if p.exists():
                files.append(p)

    qa = validate_bundle(article, out_dir, terms)
    manifest = {
        "source": str(article.path),
        "title": article.title,
        "provider": "soia-local",
        "content_modes": ["preserve", "learning", "visual_dense"],
        "files": [str(p) for p in files],
        "render": render_result,
        "qa": qa,
    }
    manifest_path = out_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    if args.strict and not qa.get("ok"):
        raise SystemExit(json.dumps(qa, ensure_ascii=False, indent=2))
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--article", required=True, help="Markdown article path")
    parser.add_argument("--out-dir", required=True, help="Output directory")
    parser.add_argument("--node-bin", default="node", help="Node.js executable for Playwright rendering")
    parser.add_argument("--node-path", default=None, help="NODE_PATH containing playwright module")
    parser.add_argument("--no-render", action="store_true", help="Skip Playwright PDF/PNG rendering")
    parser.add_argument("--strict", action="store_true", help="Exit non-zero if quality gates fail")
    parser.add_argument("--json", action="store_true", help="Print manifest JSON")
    args = parser.parse_args()
    manifest = build_bundle(args)
    if args.json:
        print(json.dumps(manifest, ensure_ascii=False, indent=2))
    else:
        print(f"wrote {len(manifest['files'])} files to {args.out_dir}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
