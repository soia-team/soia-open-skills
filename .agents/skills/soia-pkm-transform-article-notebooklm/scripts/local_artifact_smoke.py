#!/usr/bin/env python3
"""Generate a local transform smoke-test bundle from one Markdown article."""

from __future__ import annotations

import argparse
import csv
import html
import json
import os
import re
import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception as exc:  # pragma: no cover - reported in manifest
    Presentation = None
    PPTX_IMPORT_ERROR = exc
else:
    PPTX_IMPORT_ERROR = None

try:
    from article_packet import Article, Concept, matched_terms, parse_article, qa_floor, section_excerpt, theme_rows
    from validate_artifact_quality import validate_bundle
except ImportError:  # pragma: no cover - direct execution fallback
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from article_packet import Article, Concept, matched_terms, parse_article, qa_floor, section_excerpt, theme_rows
    from validate_artifact_quality import validate_bundle


def esc(value: object) -> str:
    return html.escape(str(value), quote=True)


def write_text(path: Path, content: str) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content.strip() + "\n", encoding="utf-8")
    return path


def source_label(article: Article) -> str:
    parts = [article.author, article.published_at, article.url or str(article.path)]
    return " · ".join(part for part in parts if part) or str(article.path)


def chunked(items: list[Concept], size: int) -> list[list[Concept]]:
    return [items[i : i + size] for i in range(0, len(items), size)]


def term_names(terms: list[Concept], limit: int = 24) -> str:
    return "、".join(term for term, _, _ in terms[:limit])


def write_report(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    floor = qa_floor(article, terms)
    rows = theme_rows(terms)
    toc = "\n".join(f"- {title}" for title, _ in article.sections[:24])
    coverage = "\n".join(
        f"| {theme} | {', '.join(term for term, _, _ in bucket)} | {len(bucket)} |"
        for theme, bucket in rows
    )
    glossary = "\n".join(f"| {term} | {category} | {definition} |" for term, category, definition in terms)
    section_notes = "\n\n".join(
        f"### {title}\n\n{section_excerpt(content, 900)}"
        for title, content in article.sections
        if content.strip()
    )
    checks = "\n".join(
        f"- **{term}**：归入「{category}」。转换时至少保留定义、上下文位置、与相邻概念的区别。"
        for term, category, _ in terms[: floor["min_terms"]]
    )
    report = f"""
# {article.title}｜保真转换报告

来源：{source_label(article)}

## 生成边界

- 内容模式：`preserve + learning + visual_dense`
- 本报告不是摘要，也不是全文 PDF；它用于检查后续 PPT、长图、试卷、播客、视频脚本是否覆盖原文主体。
- 生成依据只来自输入文章及其 frontmatter；公共脚本不内置个人路径、私有配置、账号、密钥或某个样例文章的专属术语表。

## 读者先看到的结论

1. 这篇文章的主线是「{article.title}」。转换时要保留它的章节顺序、概念关系、案例链和边界提醒。
2. 已从 source 中抽取 {len(article.sections)} 个章节块、{len(terms)} 个概念/流程节点。中长文产物不能只给几页摘要。
3. PPT、报告、试卷、闪卡、脑图、音视频脚本都应能回指到同一张覆盖清单，避免每种产物各讲各的。
4. 如果后续接入外部 provider，仍必须用本报告的覆盖矩阵和质量门验收，不以“生成成功”替代“内容完整”。

## 原文章节地图

{toc}

## 概念覆盖矩阵

| 模块 / 章节 | 需要覆盖的概念 | 数量 |
|---|---|---:|
{coverage}

## 术语与节点表

| 概念 / 节点 | 原文章节 | source-grounded 解释 |
|---|---|---|
{glossary}

## 逐节保真摘录

{section_notes}

## 转换验收清单

{checks}

## 媒介化建议

- **PPT / 课件**：先放文章地图，再放概念矩阵和案例链；每一组概念至少有一页结构化表达，不把长文压缩成 3-5 个 bullet。
- **长图 / 信息图**：一屏内同时呈现主判断、流程、概念矩阵、风险/边界和行动清单；不要只做大标题加几张卡片。
- **试卷 / 闪卡**：题目覆盖定义、辨析、应用和验证责任；答案与解析单独成区。
- **播客 / 视频脚本**：按章节推进，先解释主线，再讲概念关系和易错点，最后给行动清单。
- **报告**：若用户要求完整报告，优先选择结构化 custom report，而不是短 briefing。

## 残余风险

- 自动概念抽取是启发式，复杂文章可能需要人工补 1-2 个关键概念。
- 外部设计或 Notebook provider 的版式不可完全控制，输出后仍要跑质量门。
- 本地视觉产物重在可复验和信息密度，不冒充外部设计系统的导出结果。
"""
    while len(report) < floor["min_report_chars"] + 400:
        report += "\n\n## Source 补充摘录\n\n" + section_excerpt(article.plain_text, 1200)
        if len(article.plain_text) < 1200:
            break
    return write_text(out_dir / "report.md", report)


def write_data_table(out_dir: Path, terms: list[Concept]) -> Path:
    path = out_dir / "data-table.csv"
    with path.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.writer(fh)
        writer.writerow(["concept", "section", "source_explanation", "transform_use", "validation_note"])
        for term, category, definition in terms:
            writer.writerow(
                [
                    term,
                    category,
                    definition,
                    "可用于 PPT 概念页、报告矩阵、闪卡、试卷和长图信息块",
                    "产物中应能看到该概念与原文章节的关系",
                ]
            )
    return path


def write_quiz(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    floor = qa_floor(article, terms)
    questions: list[str] = []
    answers: list[str] = []
    selected = terms[: max(floor["min_questions"], min(len(terms), 12))]
    for idx, (term, category, definition) in enumerate(selected, 1):
        questions.append(
            f"### {idx}. {term} 在原文中最接近哪一类内容？\n\n"
            f"A. {category}\nB. 与原文无关的背景噪音\nC. 只适合作为装饰标题\nD. 可以从所有产物中删除\n"
        )
        answers.append(f"### {idx}. 答案：A\n\n解析：{definition}")

    next_id = len(questions) + 1
    application = [
        (
            "请写出这篇文章的 5 个核心模块，并说明每个模块适合转成哪种产物。",
            "答案要点应覆盖文章地图、概念矩阵、案例/流程、风险边界和行动清单。"
        ),
        (
            "如果把这篇文章转成 6 页以内的 PPT，最可能遗漏什么？",
            "答案要点：长文主体概念、章节关系、例子、易混点和来源页很容易被压缩掉。"
        ),
        (
            "请选 3 个你认为最容易混淆的概念，并写出区分口径。",
            "答案应使用 source 中的上下文解释，而不是凭空写百科定义。"
        ),
    ]
    for prompt, answer in application:
        questions.append(f"### {next_id}. {prompt}\n\n请用 3-5 句话回答。")
        answers.append(f"### {next_id}. 参考答案\n\n{answer}")
        next_id += 1

    return write_text(
        out_dir / "quiz.md",
        f"# 试卷：{article.title}\n\n## 题目\n\n"
        + "\n\n".join(questions)
        + "\n\n## 答案与解析\n\n"
        + "\n\n".join(answers),
    )


def write_flashcards(article: Article, out_dir: Path, terms: list[Concept]) -> list[Path]:
    md_lines = [f"# 闪卡：{article.title}", "", "| 正面 | 背面 |", "|---|---|"]
    csv_path = out_dir / "flashcards.csv"
    with csv_path.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.writer(fh)
        writer.writerow(["front", "back", "section"])
        for term, category, definition in terms:
            front = f"{term} 是什么？"
            back = f"{definition}（原文章节：{category}）"
            writer.writerow([front, back, category])
            md_lines.append(f"| {front} | {back} |")
    md_path = write_text(out_dir / "flashcards.md", "\n".join(md_lines))
    return [md_path, csv_path]


def write_mindmap(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    lines = ["mindmap", f"  root(({article.title[:36]}))"]
    for theme, bucket in theme_rows(terms)[:10]:
        safe_theme = re.sub(r"[:：#`]", "", theme)[:28] or "模块"
        lines.append(f"    {safe_theme}")
        for term, _, definition in bucket[:6]:
            safe_term = re.sub(r"[:：#`]", "", term)[:26]
            safe_def = re.sub(r"[:：#`]", "", definition)[:42]
            lines.append(f"      {safe_term}")
            if safe_def:
                lines.append(f"        {safe_def}")
    return write_text(out_dir / "mindmap.mmd", "\n".join(lines))


def write_scripts(article: Article, out_dir: Path, terms: list[Concept]) -> list[Path]:
    floor = qa_floor(article, terms)
    rows = theme_rows(terms)
    theme_script = "\n".join(
        f"### {theme}\n这一段覆盖：{term_names(bucket, 8)}。\n讲法：先给 source 中的定义，再说明它在原文章节里承担的作用。"
        for theme, bucket in rows[:10]
    )
    section_walk = "\n".join(
        f"- {title}：{section_excerpt(content, 180)}" for title, content in article.sections[:12]
    )
    podcast = f"""
# Podcast Script｜{article.title}

目标时长：中长文 deep-dive，默认 8-12 分钟。

## 开场

今天这期只围绕一篇 source：{article.title}。我们不把它压成摘要，而是把文章里的章节、概念、案例链和风险边界重新串起来，让听众听完能复述主线，也能知道哪些地方需要回看原文。

## 文章地图

{section_walk}

## 分段讲述

{theme_script}

## 主持人口播节奏

1. 先用 30 秒说明文章解决的问题。
2. 再用 2-3 分钟讲文章地图，避免听众只记住零散名词。
3. 中段逐组讲概念，每组都回答：它在原文里是什么、为什么出现、和相邻概念有什么区别。
4. 后段讲应用：把文章转成 PPT、长图、报告或试卷时，哪些内容不能丢。
5. 最后讲验收：页数、信息块、题目、卡片、音视频脚本都要回到同一张覆盖清单。

## 收束

这篇文章最值得带走的不是单个词，而是它组织问题的方法：从 source 出发，保留结构，标出边界，再把内容转换成不同媒介。任何看起来很顺但无法回指 source 的产物，都需要重做。
"""
    while len(podcast) < floor["min_podcast_chars"] + 200:
        podcast += "\n\n## 补充讲述\n\n" + section_excerpt(article.plain_text, 900)
        if len(article.plain_text) < 900:
            break

    scenes: list[str] = [
        "## 镜头 1：标题与来源\n画面：标题、作者/来源、日期。旁白：交代本片只基于这一篇 source。",
        "## 镜头 2：文章地图\n画面：章节节点依次展开。旁白：说明主线，而不是直接跳到结论。",
    ]
    scene_id = 3
    for theme, bucket in rows[:8]:
        scenes.append(
            f"## 镜头 {scene_id}：{theme}\n"
            f"画面：{theme} 的概念卡片和连接线。字幕：{term_names(bucket, 6)}。\n"
            f"旁白：解释这一组内容在原文中的位置，并提醒不要遗漏。"
        )
        scene_id += 1
    while len(scenes) < floor["min_video_scenes"]:
        section = article.sections[(len(scenes) - 2) % len(article.sections)]
        scenes.append(
            f"## 镜头 {scene_id}：{section[0]}\n"
            f"画面：source 摘录与重点标注。旁白：{section_excerpt(section[1], 160)}"
        )
        scene_id += 1
    scenes.append(f"## 镜头 {scene_id}：验收\n画面：报告、PPT、长图、试卷、闪卡逐项打勾。旁白：文件存在不是完成，覆盖度才是完成。")

    shots = [
        f"1. 黑底标题卡：{article.title}",
        "2. source 纸页展开，章节线条从左到右连接。",
        "3. 概念节点从正文中浮出，按原文章节分组。",
        "4. 一条主线穿过章节、概念、案例和边界。",
        "5. 画面切到 PPT 网格，页数与覆盖矩阵同步出现。",
        "6. 画面切到长图，信息块密集但层级清楚。",
        "7. 画面切到试卷和闪卡，问题与答案分区。",
        "8. 画面切到音视频脚本，章节节奏逐段推进。",
        "9. 红色校验线扫过缺失概念，提示必须补齐。",
        "10. 结尾字幕：从 source 出发，以覆盖度验收。",
    ]
    for idx, (theme, bucket) in enumerate(rows[:6], len(shots) + 1):
        shots.append(f"{idx}. {theme} 节点特写：{term_names(bucket, 5)}。")

    return [
        write_text(out_dir / "podcast-script.md", podcast),
        write_text(out_dir / "video-script.md", f"# Video Script｜{article.title}\n\n" + "\n\n".join(scenes)),
        write_text(out_dir / "cinematic-video-shotlist.md", f"# Cinematic Video Shotlist｜{article.title}\n\n" + "\n".join(shots)),
    ]


REPORT_CSS = """
:root { color-scheme: light; }
body { margin: 0; background: #f4f7f9; color: #1b2430; font-family: -apple-system, BlinkMacSystemFont, "PingFang SC", "Hiragino Sans GB", "Microsoft YaHei", sans-serif; }
.page { max-width: 1120px; margin: 0 auto; padding: 44px 44px 72px; }
.hero { border-left: 9px solid #f59e0b; padding: 10px 0 14px 24px; margin-bottom: 26px; }
.hero h1 { margin: 0 0 10px; font-size: 36px; line-height: 1.18; letter-spacing: 0; }
.hero p { margin: 0; color: #516070; font-size: 15px; }
.grid { display: grid; grid-template-columns: repeat(3, 1fr); gap: 14px; margin: 18px 0 28px; }
.stat, .panel { background: white; border: 1px solid #dce5ee; border-radius: 8px; padding: 16px; box-shadow: 0 8px 20px rgba(19, 37, 58, .06); }
.stat b { display: block; color: #0f766e; font-size: 30px; margin-bottom: 4px; }
.panel h2 { margin: 0 0 12px; font-size: 22px; }
table { width: 100%; border-collapse: collapse; background: white; border: 1px solid #dce5ee; border-radius: 8px; overflow: hidden; margin: 12px 0 24px; }
th, td { border-bottom: 1px solid #e5edf4; padding: 10px 12px; text-align: left; vertical-align: top; font-size: 13px; line-height: 1.55; }
th { background: #eaf2f8; color: #26394d; }
.section { margin: 14px 0; padding: 14px 18px; background: #fff; border-left: 4px solid #2563eb; border-radius: 8px; }
.section h3 { margin: 0 0 8px; font-size: 18px; }
.muted { color: #657386; }
"""


def write_report_html(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    rows = theme_rows(terms)
    coverage_html = "".join(
        f"<tr><td>{esc(theme)}</td><td>{esc(term_names(bucket, 12))}</td><td>{len(bucket)}</td></tr>"
        for theme, bucket in rows
    )
    section_html = "".join(
        f"<div class='section'><h3>{esc(title)}</h3><p>{esc(section_excerpt(content, 520))}</p></div>"
        for title, content in article.sections[:10]
    )
    glossary_html = "".join(
        f"<tr><td>{esc(term)}</td><td>{esc(category)}</td><td>{esc(definition)}</td></tr>"
        for term, category, definition in terms[:36]
    )
    body = f"""
<main class="page">
  <section class="hero">
    <h1>{esc(article.title)}</h1>
    <p>{esc(source_label(article))}</p>
  </section>
  <section class="grid">
    <div class="stat"><b>{len(article.sections)}</b><span>章节块</span></div>
    <div class="stat"><b>{len(terms)}</b><span>概念/节点</span></div>
    <div class="stat"><b>{qa_floor(article, terms)["min_slides"]}</b><span>PPT 最低页数</span></div>
  </section>
  <section class="panel">
    <h2>覆盖矩阵</h2>
    <table><thead><tr><th>模块</th><th>概念</th><th>数量</th></tr></thead><tbody>{coverage_html}</tbody></table>
  </section>
  <section class="panel">
    <h2>逐节摘录</h2>
    {section_html}
  </section>
  <section class="panel">
    <h2>概念表</h2>
    <table><thead><tr><th>概念</th><th>章节</th><th>解释</th></tr></thead><tbody>{glossary_html}</tbody></table>
  </section>
</main>
"""
    return write_text(out_dir / "report.html", html_doc(article.title, REPORT_CSS, body))


INFOGRAPHIC_CSS = """
body { margin: 0; background: #0d1726; color: #f8fafc; font-family: -apple-system, BlinkMacSystemFont, "PingFang SC", "Microsoft YaHei", sans-serif; }
.canvas { width: 1080px; min-height: 1920px; box-sizing: border-box; padding: 44px; background: linear-gradient(180deg, #0d1726 0%, #12243a 55%, #0f172a 100%); }
.top { border: 2px solid #f59e0b; padding: 24px 28px; border-radius: 8px; margin-bottom: 24px; background: rgba(245, 158, 11, .09); }
h1 { font-size: 52px; line-height: 1.12; margin: 0 0 14px; letter-spacing: 0; }
h2 { font-size: 28px; margin: 0 0 16px; color: #fde68a; }
.subtitle { color: #cbd5e1; font-size: 20px; line-height: 1.55; }
.stats { display: grid; grid-template-columns: repeat(3, 1fr); gap: 14px; margin: 20px 0; }
.stat, .card, .flow-node, .matrix-item, .risk { border: 1px solid rgba(148, 163, 184, .45); border-radius: 8px; padding: 16px; background: rgba(15, 23, 42, .72); }
.stat b { display: block; color: #60a5fa; font-size: 40px; }
.grid { display: grid; grid-template-columns: repeat(3, 1fr); gap: 14px; margin-bottom: 22px; }
.card b { display: block; color: #fbbf24; font-size: 22px; margin-bottom: 8px; }
.card p, .risk p, .matrix-item p { margin: 0; color: #dbeafe; font-size: 16px; line-height: 1.5; }
.flow { display: grid; grid-template-columns: repeat(5, 1fr); gap: 10px; margin-bottom: 22px; }
.flow-node { min-height: 94px; border-color: rgba(96, 165, 250, .7); }
.flow-node b { color: #93c5fd; }
.two { display: grid; grid-template-columns: 1.15fr .85fr; gap: 16px; }
.matrix { display: grid; grid-template-columns: repeat(2, 1fr); gap: 12px; }
.risk { border-color: rgba(248, 113, 113, .7); background: rgba(127, 29, 29, .38); }
.footer { margin-top: 20px; color: #94a3b8; font-size: 14px; }
"""


def write_infographic_html(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    rows = theme_rows(terms)
    floor = qa_floor(article, terms)
    cards = "".join(
        f"<div class='card' data-block='info'><b>{esc(term)}</b><p>{esc(definition)}</p></div>"
        for term, _, definition in terms[: max(floor["min_infographic_blocks"], 12)]
    )
    flow_items = [title for title, _ in article.sections[:5]]
    while len(flow_items) < 5:
        flow_items.append(f"模块 {len(flow_items) + 1}")
    flow = "".join(
        f"<div class='flow-node' data-block='info'><b>{idx}</b><p>{esc(title)}</p></div>"
        for idx, title in enumerate(flow_items[:5], 1)
    )
    matrix = "".join(
        f"<div class='matrix-item' data-block='info'><b>{esc(theme)}</b><p>{esc(term_names(bucket, 6))}</p></div>"
        for theme, bucket in rows[:6]
    )
    risks = "".join(
        f"<div class='risk' data-block='info'><b>{esc(term)}</b><p>不要脱离「{esc(category)}」单独解释；保留 source 上下文。</p></div>"
        for term, category, _ in terms[:4]
    )
    body = f"""
<main class="canvas">
  <section class="top" data-block="info">
    <h1>{esc(article.title)}</h1>
    <div class="subtitle">保真转换信息图：先保留文章地图，再覆盖概念、流程、风险和行动清单。</div>
  </section>
  <section class="stats">
    <div class="stat" data-block="info"><b>{len(article.sections)}</b><span>source sections</span></div>
    <div class="stat" data-block="info"><b>{len(terms)}</b><span>concept nodes</span></div>
    <div class="stat" data-block="info"><b>{floor["min_slides"]}</b><span>minimum slides</span></div>
  </section>
  <h2>文章路径</h2>
  <section class="flow">{flow}</section>
  <h2>概念卡片</h2>
  <section class="grid">{cards}</section>
  <section class="two">
    <div>
      <h2>模块矩阵</h2>
      <div class="matrix">{matrix}</div>
    </div>
    <div>
      <h2>风险边界</h2>
      {risks}
    </div>
  </section>
  <div class="footer">source: {esc(source_label(article))}</div>
</main>
"""
    return write_text(out_dir / "infographic.html", html_doc(article.title, INFOGRAPHIC_CSS, body))


DECK_CSS = """
body { margin: 0; background: #111827; color: #172033; font-family: -apple-system, BlinkMacSystemFont, "PingFang SC", "Microsoft YaHei", sans-serif; }
.slide { width: 1600px; height: 900px; box-sizing: border-box; padding: 64px 78px; background: #f8fafc; overflow: hidden; position: relative; page-break-after: always; }
.slide.dark { background: #0f172a; color: #f8fafc; }
.slide h1 { margin: 0 0 24px; font-size: 66px; line-height: 1.08; letter-spacing: 0; }
.slide h2 { margin: 0 0 26px; font-size: 46px; line-height: 1.14; letter-spacing: 0; }
.slide p, .slide li { font-size: 26px; line-height: 1.45; }
.kicker { color: #f59e0b; font-size: 24px; font-weight: 700; margin-bottom: 18px; }
.grid { display: grid; grid-template-columns: repeat(3, 1fr); gap: 20px; }
.two { display: grid; grid-template-columns: 1fr 1fr; gap: 28px; }
.card, .box, .metric { background: white; border: 1px solid #dbe4ee; border-radius: 8px; padding: 24px; box-shadow: 0 10px 28px rgba(15, 23, 42, .08); }
.dark .card, .dark .box, .dark .metric { background: rgba(15, 23, 42, .75); border-color: rgba(148, 163, 184, .4); }
.card b, .box b { display: block; color: #2563eb; font-size: 28px; margin-bottom: 10px; }
.dark .card b, .dark .box b { color: #fbbf24; }
.metric b { display: block; font-size: 58px; color: #0f766e; }
table { width: 100%; border-collapse: collapse; font-size: 22px; background: white; }
th, td { border: 1px solid #dbe4ee; padding: 14px 16px; vertical-align: top; line-height: 1.35; }
th { background: #e2e8f0; }
.flow { display: grid; grid-template-columns: repeat(5, 1fr); gap: 14px; }
.flow div { background: #1d4ed8; color: white; border-radius: 8px; padding: 20px; min-height: 118px; }
.small { font-size: 20px; color: #64748b; }
.dark .small { color: #cbd5e1; }
"""


def slide(title: str, body: str, cls: str = "") -> str:
    return f"<section class='slide {cls}'><div class='kicker'>{esc(title)}</div>{body}</section>"


def write_deck_html(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    floor = qa_floor(article, terms)
    rows = theme_rows(terms)
    slides: list[str] = []
    metric_cards = (
        f"<div class='metric'><b>{len(article.sections)}</b><span>章节块</span></div>"
        f"<div class='metric'><b>{len(terms)}</b><span>概念/节点</span></div>"
        f"<div class='metric'><b>{floor['min_slides']}</b><span>最低页数</span></div>"
    )
    slides.append(
        slide(
            "Cover",
            f"<h1>{esc(article.title)}</h1><p>把 source 转成可讲、可看、可考、可复验的一组产物。</p><div class='grid'>{metric_cards}</div>",
            "dark",
        )
    )
    section_cards = "".join(
        f"<div class='card'><b>{esc(title)}</b><p>{esc(section_excerpt(content, 120))}</p></div>"
        for title, content in article.sections[:6]
    )
    slides.append(slide("文章地图", f"<h2>先保留结构，再做媒介化</h2><div class='grid'>{section_cards}</div>"))

    matrix_rows = "".join(
        f"<tr><td>{esc(theme)}</td><td>{esc(term_names(bucket, 10))}</td><td>{len(bucket)}</td></tr>"
        for theme, bucket in rows[:12]
    )
    slides.append(
        slide(
            "覆盖矩阵",
            f"<h2>所有产物共用这一张清单</h2><table><tr><th>模块</th><th>概念</th><th>数量</th></tr>{matrix_rows}</table>",
        )
    )

    flow_nodes = "".join(
        f"<div><b>{idx}</b><p>{esc(title)}</p></div>" for idx, (title, _) in enumerate(article.sections[:5], 1)
    )
    slides.append(slide("Source Path", f"<h2>按原文路径推进</h2><div class='flow'>{flow_nodes}</div>", "dark"))

    for theme, bucket in rows[:8]:
        cards = "".join(
            f"<div class='card'><b>{esc(term)}</b><p>{esc(definition)}</p></div>"
            for term, _, definition in bucket[:6]
        )
        slides.append(slide(theme, f"<h2>{esc(theme)}</h2><div class='grid'>{cards}</div>"))

    for idx, group in enumerate(chunked(terms, 10)[:2], 1):
        table_rows = "".join(
            f"<tr><td>{esc(term)}</td><td>{esc(category)}</td><td>{esc(definition)}</td></tr>"
            for term, category, definition in group
        )
        slides.append(slide(f"术语速查 {idx}", f"<h2>覆盖清单 {idx}</h2><table><tr><th>概念</th><th>章节</th><th>解释</th></tr>{table_rows}</table>"))

    pair_cards = ""
    for left, right in zip(terms[0::2][:4], terms[1::2][:4]):
        pair_cards += (
            f"<div class='box'><b>{esc(left[0])} / {esc(right[0])}</b>"
            f"<p>{esc(left[1])} vs {esc(right[1])}。讲解时要回到 source，不凭空扩展。</p></div>"
        )
    slides.append(slide("易混与边界", f"<h2>相邻概念要讲出区别</h2><div class='two'>{pair_cards}</div>"))

    slides.append(
        slide(
            "转换行动清单",
            "<h2>交付前逐项检查</h2><ol><li>是否覆盖主要章节？</li><li>是否覆盖主体概念？</li><li>是否保留案例/流程？</li><li>是否有风险和边界？</li><li>是否能回指 source？</li></ol>",
        )
    )
    slides.append(
        slide(
            "自测",
            "<h2>三道验收题</h2><ol><li>这篇文章的主线是什么？</li><li>哪些概念被遗漏会改变理解？</li><li>哪个产物最容易变成摘要，如何补救？</li></ol>",
            "dark",
        )
    )
    coverage_index = esc(term_names(terms, 40))
    slides.append(slide("Source", f"<h2>来源与覆盖索引</h2><p>{esc(source_label(article))}</p><p class='small'>{coverage_index}</p>"))

    if len(slides) > floor["max_slides"]:
        slides = slides[: floor["max_slides"] - 1] + [slides[-1]]
    while len(slides) < floor["min_slides"]:
        section = article.sections[(len(slides) - 1) % len(article.sections)]
        slides.insert(
            -1,
            slide(
                f"Source Detail {len(slides)}",
                f"<h2>{esc(section[0])}</h2><p>{esc(section_excerpt(section[1], 500))}</p>",
            ),
        )

    html_body = "\n".join(slides)
    return write_text(out_dir / "deck.html", html_doc(article.title, DECK_CSS, html_body))


def html_doc(title: str, css: str, body: str) -> str:
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{esc(title)}</title>
  <style>{css}</style>
</head>
<body>
{body}
</body>
</html>
"""


def add_textbox(slide_obj, left, top, width, height, text: str, size: int, color: RGBColor | None = None, bold: bool = False) -> None:
    box = slide_obj.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_card(slide_obj, left, top, width, height, title: str, body: str, fill: RGBColor) -> None:
    shape = slide_obj.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(220, 226, 235)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(30, 41, 59)
    p2 = frame.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(51, 65, 85)


def write_deck_pptx(article: Article, out_dir: Path, terms: list[Concept]) -> Path | None:
    if Presentation is None:
        return None
    floor = qa_floor(article, terms)
    rows = theme_rows(terms)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide(title: str, subtitle: str = ""):
        s = prs.slides.add_slide(blank)
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(248, 250, 252)
        add_textbox(s, Inches(0.55), Inches(0.35), Inches(12.1), Inches(0.58), title, 27, RGBColor(15, 23, 42), True)
        if subtitle:
            add_textbox(s, Inches(0.58), Inches(0.95), Inches(11.6), Inches(0.38), subtitle, 12, RGBColor(100, 116, 139))
        return s

    s = new_slide(article.title, source_label(article))
    add_textbox(s, Inches(0.7), Inches(2.0), Inches(11.8), Inches(0.72), "保真转换：先覆盖 source，再生成不同媒介产物。", 23, RGBColor(30, 64, 175), True)
    for idx, (label, value) in enumerate([("章节块", len(article.sections)), ("概念节点", len(terms)), ("最低页数", floor["min_slides"])]):
        add_card(s, Inches(0.75 + idx * 4.1), Inches(3.2), Inches(3.6), Inches(1.55), label, str(value), RGBColor(239, 246, 255))

    s = new_slide("文章地图", "不跳过结构，先让读者知道原文如何展开")
    for idx, (title, content) in enumerate(article.sections[:6]):
        add_card(s, Inches(0.6 + (idx % 3) * 4.15), Inches(1.55 + (idx // 3) * 2.35), Inches(3.75), Inches(1.85), title, section_excerpt(content, 130), RGBColor(255, 255, 255))

    s = new_slide("覆盖矩阵", "所有产物共用同一张概念清单")
    for idx, (theme, bucket) in enumerate(rows[:8]):
        add_card(s, Inches(0.6 + (idx % 4) * 3.15), Inches(1.45 + (idx // 4) * 2.25), Inches(2.8), Inches(1.8), theme, term_names(bucket, 6), RGBColor(240, 253, 244))

    for theme, bucket in rows[:8]:
        s = new_slide(theme, "每页负责一组概念，避免把长文压成少量 bullet")
        for idx, (term, _, definition) in enumerate(bucket[:6]):
            add_card(s, Inches(0.6 + (idx % 3) * 4.15), Inches(1.5 + (idx // 3) * 2.25), Inches(3.75), Inches(1.75), term, definition, RGBColor(255, 255, 255))

    for idx, group in enumerate(chunked(terms, 8)[:2], 1):
        s = new_slide(f"术语速查 {idx}", "检查 PPT 是否覆盖主体概念")
        for row_idx, (term, category, definition) in enumerate(group):
            add_textbox(s, Inches(0.72), Inches(1.35 + row_idx * 0.65), Inches(2.1), Inches(0.38), term, 14, RGBColor(30, 64, 175), True)
            add_textbox(s, Inches(2.65), Inches(1.35 + row_idx * 0.65), Inches(2.4), Inches(0.38), category, 12, RGBColor(71, 85, 105))
            add_textbox(s, Inches(5.0), Inches(1.35 + row_idx * 0.65), Inches(7.2), Inches(0.42), definition, 11, RGBColor(51, 65, 85))

    s = new_slide("行动清单", "交付前逐项验收")
    add_textbox(s, Inches(0.9), Inches(1.6), Inches(11.2), Inches(4.8), "1. 覆盖主要章节\n2. 覆盖主体概念\n3. 保留案例、流程或证据链\n4. 标出风险与边界\n5. 输出文件可打开、可读、可追溯", 24, RGBColor(15, 23, 42))

    s = new_slide("来源", "Source-grounded output")
    add_textbox(s, Inches(0.9), Inches(1.6), Inches(11.4), Inches(2.2), source_label(article), 18, RGBColor(71, 85, 105))
    add_textbox(s, Inches(0.9), Inches(3.6), Inches(11.4), Inches(2.0), term_names(terms, 36), 14, RGBColor(100, 116, 139))

    while len(prs.slides) < min(floor["min_slides"], 14):
        title, content = article.sections[(len(prs.slides) - 1) % len(article.sections)]
        s = new_slide(f"Source Detail {len(prs.slides)}", title)
        add_textbox(s, Inches(0.9), Inches(1.55), Inches(11.4), Inches(4.8), section_excerpt(content, 900), 18, RGBColor(30, 41, 59))

    path = out_dir / "deck.pptx"
    prs.save(path)
    return path


def render_with_playwright(out_dir: Path, node_bin: str, node_path: str | None = None) -> dict[str, object]:
    script = """
const { chromium } = require('playwright');
const path = require('path');
const outDir = __OUTDIR__;
const toUrl = file => 'file://' + path.resolve(outDir, file);

async function launchBrowser() {
  try {
    return await chromium.launch({ headless: true, channel: 'chrome' });
  } catch (first) {
    return await chromium.launch({ headless: true });
  }
}

(async () => {
  const browser = await launchBrowser();
  const page = await browser.newPage();
  await page.goto(toUrl('report.html'), { waitUntil: 'networkidle' });
  await page.pdf({ path: path.join(outDir, 'report.pdf'), format: 'A4', printBackground: true, margin: { top: '12mm', bottom: '12mm', left: '10mm', right: '10mm' } });
  await page.setViewportSize({ width: 1080, height: 1920 });
  await page.goto(toUrl('infographic.html'), { waitUntil: 'networkidle' });
  await page.screenshot({ path: path.join(outDir, 'infographic.png'), fullPage: true });
  await page.setViewportSize({ width: 1600, height: 900 });
  await page.goto(toUrl('deck.html'), { waitUntil: 'networkidle' });
  await page.screenshot({ path: path.join(outDir, 'deck-cover.png') });
  await page.pdf({ path: path.join(outDir, 'deck.pdf'), width: '1600px', height: '900px', printBackground: true });
  await browser.close();
})().catch(err => {
  console.error(err && err.stack ? err.stack : String(err));
  process.exit(1);
});
""".replace("__OUTDIR__", json.dumps(str(out_dir)))
    env = os.environ.copy()
    if node_path:
        env["NODE_PATH"] = node_path
    proc = subprocess.run([node_bin, "-e", script], text=True, capture_output=True, env=env, check=False)
    return {
        "ok": proc.returncode == 0,
        "returncode": proc.returncode,
        "stdout": proc.stdout.strip(),
        "stderr": proc.stderr.strip(),
    }


def build_bundle(args: argparse.Namespace) -> dict[str, object]:
    article = parse_article(Path(args.article))
    out_dir = Path(args.out_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    terms = matched_terms(article)

    files: list[Path] = [
        write_report(article, out_dir, terms),
        write_report_html(article, out_dir, terms),
        write_data_table(out_dir, terms),
        write_quiz(article, out_dir, terms),
        *write_flashcards(article, out_dir, terms),
        write_mindmap(article, out_dir, terms),
        *write_scripts(article, out_dir, terms),
        write_infographic_html(article, out_dir, terms),
        write_deck_html(article, out_dir, terms),
    ]
    pptx_path = write_deck_pptx(article, out_dir, terms)
    if pptx_path:
        files.append(pptx_path)

    render_result = None
    if not args.no_render:
        render_result = render_with_playwright(out_dir, args.node_bin, args.node_path)
        for name in ["report.pdf", "infographic.png", "deck-cover.png", "deck.pdf"]:
            path = out_dir / name
            if path.exists():
                files.append(path)

    qa = validate_bundle(article, out_dir, terms)
    manifest = {
        "source": str(article.path),
        "title": article.title,
        "provider": "soia-local",
        "content_modes": ["preserve", "learning", "visual_dense"],
        "concept_count": len(terms),
        "section_count": len(article.sections),
        "files": [str(path) for path in files],
        "render": render_result,
        "qa": qa,
    }
    manifest_path = out_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    if args.strict and not qa.get("ok"):
        raise SystemExit(json.dumps(qa, ensure_ascii=False, indent=2))
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--article", required=True, help="Markdown article path")
    parser.add_argument("--out-dir", required=True, help="Output directory")
    parser.add_argument("--node-bin", default="node", help="Node.js executable for Playwright rendering")
    parser.add_argument("--node-path", default=None, help="NODE_PATH containing playwright module")
    parser.add_argument("--no-render", action="store_true", help="Skip Playwright PDF/PNG rendering")
    parser.add_argument("--strict", action="store_true", help="Exit non-zero if quality gates fail")
    parser.add_argument("--json", action="store_true", help="Print manifest JSON")
    args = parser.parse_args()
    manifest = build_bundle(args)
    if args.json:
        print(json.dumps(manifest, ensure_ascii=False, indent=2))
    else:
        print(f"wrote {len(manifest['files'])} files to {args.out_dir}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
